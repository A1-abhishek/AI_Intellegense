import io
import json
import csv
import logging
import os
import re
import struct
import subprocess
import tempfile
import zipfile
import xml.etree.ElementTree as ET
from pathlib import Path

logger = logging.getLogger("docmind.doc_processor")


def extract_text_from_bytes(data: bytes, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    logger.debug(f"Extracting text from {filename} (ext={ext}, size={len(data)})")

    extractors = {
        ".txt": _plain, ".md": _plain, ".log": _plain,
        ".csv": _csv_text, ".json": _json_text, ".xml": _xml_text,
        ".html": _html_text, ".htm": _html_text,
        ".yaml": _plain, ".yml": _plain,
        ".pdf": _pdf, ".docx": _docx, ".doc": _docx,
        ".pptx": _pptx, ".xlsx": _xlsx,
        ".odt": _odt, ".rtf": _rtf, ".epub": _epub,
    }

    extractor = extractors.get(ext)
    if extractor:
        try:
            result = extractor(data, filename)
            logger.info(f"Extracted {len(result)} chars from {filename}")
            return result
        except Exception as e:
            logger.warning(f"Primary extraction failed for {filename}: {e}, trying LibreOffice fallback")
            try:
                result = _libreoffice_to_pdf(data, filename)
                logger.info(f"LibreOffice fallback extracted {len(result)} chars from {filename}")
                return result
            except Exception as e2:
                logger.error(f"LibreOffice fallback also failed for {filename}: {e2}")
                raise e

    logger.warning(f"No extractor for {ext}, trying raw decode")
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")


def _plain(data: bytes, filename: str) -> str:
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1")


def _csv_text(data: bytes, filename: str) -> str:
    text = data.decode("utf-8", errors="replace")
    reader = csv.reader(io.StringIO(text))
    rows = list(reader)
    lines = []
    for row in rows:
        lines.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(lines)


def _json_text(data: bytes, filename: str) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        parsed = json.loads(text)
        return json.dumps(parsed, indent=2, ensure_ascii=False)
    except json.JSONDecodeError:
        return text


def _xml_text(data: bytes, filename: str) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        root = ET.fromstring(text)
        return _etree_to_text(root)
    except ET.ParseError:
        return text


def _etree_to_text(element, indent=0):
    lines = []
    prefix = "  " * indent
    if element.text and element.text.strip():
        lines.append(f"{prefix}{element.text.strip()}")
    for child in element:
        lines.append(f"{prefix}{child.tag}:")
        lines.append(_etree_to_text(child, indent + 1))
    if element.tail and element.tail.strip():
        lines.append(f"{prefix}{element.tail.strip()}")
    return "\n".join(lines)


def _html_text(data: bytes, filename: str) -> str:
    text = data.decode("utf-8", errors="replace")
    try:
        from bs4 import BeautifulSoup
        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return soup.get_text(separator="\n", strip=True)
    except ImportError:
        import re
        clean = re.sub(r"<[^>]+>", " ", text)
        return re.sub(r"\s+", " ", clean).strip()


def _pdf(data: bytes, filename: str) -> str:
    logger.debug(f"Extracting PDF: {filename}")
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        pages = []
        for i, page in enumerate(doc):
            text = page.get_text()
            pages.append(text)
            logger.debug(f"PDF page {i+1}: {len(text)} chars")
        doc.close()
        result = "\n\n".join(pages)
        logger.info(f"PDF extracted: {len(result)} chars from {len(pages)} pages")
        return result
    except ImportError:
        logger.error("PyMuPDF not installed")
        raise ImportError("Install PyMuPDF: pip install PyMuPDF")


def _docx(data: bytes, filename: str) -> str:
    if data[:8] == b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1":
        return _doc_ole2(data, filename)
    try:
        from docx import Document
        doc = Document(io.BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                paragraphs.append(" | ".join(cells))
        return "\n".join(paragraphs)
    except zipfile.BadZipFile:
        logger.warning(f"File {filename} has .docx extension but is not ZIP; trying OLE2 fallback")
        try:
            return _doc_ole2(data, filename)
        except Exception:
            logger.warning(f"OLE2 also failed for {filename}; trying raw text decode")
            return _plain(data, filename)
    except ImportError:
        raise ImportError("Install python-docx: pip install python-docx")


def _doc_ole2(data: bytes, filename: str) -> str:
    import re
    import struct
    try:
        import olefile
        ole = olefile.OleFileIO(io.BytesIO(data))
        if ole.exists("WordDocument"):
            word_stream = ole.openstream("WordDocument").read()
            text_pieces = []
            if ole.exists("1Table"):
                table_stream = ole.openstream("1Table").read()
            elif ole.exists("0Table"):
                table_stream = ole.openstream("0Table").read()
            else:
                table_stream = b""

            magic = struct.unpack_from("<H", word_stream, 0)[0]
            flags = struct.unpack_from("<H", word_stream, 0x0A)[0]
            nfc = struct.unpack_from("<H", word_stream, 0x0044)[0]

            if flags & 0x0004:
                fc_min = struct.unpack_from("<I", word_stream, 0x0018)[0]
                fc_max = struct.unpack_from("<I", word_stream, 0x001C)[0]
                text_bytes = word_stream[fc_min:fc_max]
            else:
                text_bytes = word_stream[0x0785:]

            decoded = text_bytes.decode("utf-16-le", errors="ignore")
            cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", decoded)
            cleaned = re.sub(r"\r?\n", "\n", cleaned)
            ole.close()
            if len(cleaned.strip()) > 10:
                logger.info(f"OLE2 text extracted from {filename}: {len(cleaned)} chars")
                return cleaned.strip()

        streams = ole.listdir()
        all_text = []
        for stream_path in streams:
            try:
                content = ole.openstream(stream_path).read()
                if len(content) > 100:
                    try:
                        text = content.decode("utf-16-le", errors="ignore")
                    except Exception:
                        text = content.decode("latin-1", errors="ignore")
                    cleaned = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
                    if len(cleaned.strip()) > 20:
                        all_text.append(cleaned.strip())
            except Exception:
                continue
        ole.close()
        if all_text:
            result = "\n\n".join(all_text)
            logger.info(f"OLE2 stream extraction from {filename}: {len(result)} chars")
            return result
    except ImportError:
        raise ImportError("Install olefile: pip install olefile")
    except Exception as e:
        logger.error(f"OLE2 extraction failed for {filename}: {e}")
    raise ValueError(f"File '{filename}' appears to be an old .doc format but text could not be extracted. Try converting to .docx first.")


def _pptx(data: bytes, filename: str) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
        return "\n".join(texts)
    except ImportError:
        raise ImportError("Install python-pptx: pip install python-pptx")


def _xlsx(data: bytes, filename: str) -> str:
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                lines.append(" | ".join(cells))
        wb.close()
        return "\n".join(lines)
    except ImportError:
        raise ImportError("Install openpyxl: pip install openpyxl")


def _odt(data: bytes, filename: str) -> str:
    try:
        from odf import text, teletype
        from odf.opendocument import load
        doc = load(io.BytesIO(data))
        all_text = []
        for paragraph in doc.spreadsheet.getElementsByType(text.P):
            all_text.append(teletype.extractText(paragraph))
        return "\n".join(all_text)
    except ImportError:
        raise ImportError("Install odfpy: pip install odfpy")


def _rtf(data: bytes, filename: str) -> str:
    try:
        from striprtf.striprtf import rtf_to_text
        text = data.decode("utf-8", errors="replace")
        return rtf_to_text(text)
    except ImportError:
        import re
        text = data.decode("utf-8", errors="replace")
        clean = re.sub(r"\{[^{}]*\}", "", text)
        clean = re.sub(r"\\[a-z]+\d*\s?", "", clean)
        clean = re.sub(r"[{}\\]", "", clean)
        return clean.strip()


def _epub(data: bytes, filename: str) -> str:
    try:
        import ebooklib
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(io.BytesIO(data))
        texts = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), "html.parser")
                texts.append(soup.get_text(separator="\n", strip=True))
        return "\n".join(texts)
    except ImportError:
        raise ImportError("Install ebooklib: pip install ebooklib")


LIBREOFFICE_PATH = os.getenv(
    "LIBREOFFICE_PATH",
    r"C:\Program Files\LibreOffice\program\soffice.exe",
)


def _libreoffice_to_pdf(data: bytes, filename: str) -> str:
    """Convert any document to PDF via LibreOffice, then extract text with PyMuPDF."""
    with tempfile.TemporaryDirectory() as tmp:
        src = os.path.join(tmp, filename)
        with open(src, "wb") as f:
            f.write(data)

        subprocess.run(
            [
                LIBREOFFICE_PATH,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmp,
                src,
            ],
            timeout=120,
            capture_output=True,
        )

        pdf_name = Path(filename).stem + ".pdf"
        pdf_path = os.path.join(tmp, pdf_name)
        if not os.path.exists(pdf_path):
            raise RuntimeError(f"LibreOffice conversion failed for {filename}")

        with open(pdf_path, "rb") as f:
            pdf_data = f.read()

        return _pdf(pdf_data, pdf_name)


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> list[dict]:
    if not text.strip():
        return []

    chunks = []
    start = 0
    idx = 0

    while start < len(text):
        end = start + chunk_size

        if end < len(text):
            for sep in ["\n\n", "\n", ". ", " "]:
                last = text.rfind(sep, start, end)
                if last > start + chunk_size // 2:
                    end = last + len(sep)
                    break

        chunk_text_val = text[start:end].strip()
        if chunk_text_val:
            chunks.append({
                "chunk_id": idx,
                "text": chunk_text_val,
                "start_char": start,
                "end_char": end,
            })
            idx += 1

        start = end - overlap if end < len(text) else end

    logger.debug(f"Chunked text into {len(chunks)} chunks")
    return chunks
